#!/usr/bin/env python3
"""
Generate Progress Report PowerPoint from Git Commits
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import subprocess
from collections import defaultdict
from datetime import datetime

def get_git_commits():
    """Get all git commits with details"""
    cmd = ['git', 'log', '--pretty=format:%H|%an|%ad|%s', '--date=short']
    result = subprocess.run(cmd, capture_output=True, text=True)

    commits = []
    for line in result.stdout.strip().split('\n'):
        if line:
            parts = line.split('|', 3)
            if len(parts) == 4:
                commits.append({
                    'hash': parts[0][:7],
                    'author': parts[1],
                    'date': parts[2],
                    'message': parts[3]
                })
    return commits

def categorize_commits(commits):
    """Categorize commits by type"""
    categories = {
        'feat': [],
        'docs': [],
        'fix': [],
        'refactor': [],
        'chore': [],
        'test': [],
        'other': []
    }

    for commit in commits:
        msg = commit['message'].lower()
        categorized = False

        for category in categories.keys():
            if msg.startswith(category + ':') or msg.startswith(category + '('):
                categories[category].append(commit)
                categorized = True
                break

        if not categorized:
            categories['other'].append(commit)

    return categories

def extract_key_features(categories):
    """Extract key features from commits"""
    features = []

    # Group by major features
    feature_groups = defaultdict(list)

    for commit in categories['feat']:
        msg = commit['message']

        # Identify major feature groups
        if 'UAV' in msg or 'uav' in msg.lower():
            feature_groups['UAV Policy xApp'].append(msg)
        elif 'Beam' in msg or 'beam' in msg.lower() or 'KPI Query' in msg:
            feature_groups['Beam KPI Query System'].append(msg)
        elif 'dual-path' in msg.lower() or 'RMR' in msg or 'HTTP' in msg:
            feature_groups['Dual-Path Communication'].append(msg)
        elif 'Grafana' in msg or 'Prometheus' in msg or 'monitoring' in msg.lower():
            feature_groups['Monitoring & Visualization'].append(msg)
        elif 'Phase' in msg:
            feature_groups['Platform Deployment'].append(msg)
        elif 'health check' in msg.lower() or 'endpoint' in msg.lower():
            feature_groups['Health Check & Endpoints'].append(msg)
        else:
            feature_groups['Other Features'].append(msg)

    return feature_groups

def add_title_slide(prs, title, subtitle):
    """Add title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])

    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]

    title_shape.text = title
    subtitle_shape.text = subtitle

    # Style title
    title_frame = title_shape.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 51, 102)

def add_section_header(prs, title):
    """Add section header slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[2])

    title_shape = slide.shapes.title
    title_shape.text = title

    # Style
    title_frame = title_shape.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 102, 204)

def add_bullet_slide(prs, title, bullets):
    """Add bullet point slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]

    title_shape.text = title

    # Style title
    title_frame = title_shape.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 51, 102)

    # Add bullets
    text_frame = body_shape.text_frame
    text_frame.clear()

    for bullet in bullets:
        p = text_frame.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(18)
        p.space_before = Pt(8)

def add_two_column_slide(prs, title, left_items, right_items):
    """Add two-column slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    # Title
    title_shape = slide.shapes.title
    title_shape.text = title

    title_frame = title_shape.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 51, 102)

    # Left column
    left = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(4.5), Inches(5))
    left_tf = left.text_frame
    left_tf.word_wrap = True

    for item in left_items:
        p = left_tf.add_paragraph()
        p.text = item
        p.font.size = Pt(14)
        p.space_before = Pt(6)

    # Right column
    right = slide.shapes.add_textbox(Inches(5.5), Inches(2), Inches(4.5), Inches(5))
    right_tf = right.text_frame
    right_tf.word_wrap = True

    for item in right_items:
        p = right_tf.add_paragraph()
        p.text = item
        p.font.size = Pt(14)
        p.space_before = Pt(6)

def generate_presentation():
    """Generate the complete presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Get commits
    print("Fetching git commits...")
    commits = get_git_commits()
    categories = categorize_commits(commits)
    feature_groups = extract_key_features(categories)

    print(f"Found {len(commits)} commits")

    # Title slide
    add_title_slide(
        prs,
        "O-RAN RIC Platform",
        f"開發進度報告\n{datetime.now().strftime('%Y-%m-%d')}"
    )

    # Overview slide
    add_section_header(prs, "專案概覽")

    overview_bullets = [
        f"總計提交次數: {len(commits)} commits",
        f"功能開發: {len(categories['feat'])} commits",
        f"文檔更新: {len(categories['docs'])} commits",
        f"錯誤修復: {len(categories['fix'])} commits",
        f"程式碼重構: {len(categories['refactor'])} commits",
        f"其他維護: {len(categories['chore'])} commits"
    ]
    add_bullet_slide(prs, "提交統計", overview_bullets)

    # Key Features Section
    add_section_header(prs, "主要功能開發")

    # UAV Policy xApp
    if 'UAV Policy xApp' in feature_groups:
        uav_features = [
            "完成 UAV Policy xApp 實作",
            "整合 TRACTOR 系統",
            "實現完整自動化流程",
            "包含 TDD 測試與文檔",
            "部署 ns-oran-bridge 整合",
            "完成端到端工作流程測試"
        ]
        add_bullet_slide(prs, "UAV Policy xApp", uav_features)

    # Beam KPI Query System
    if 'Beam KPI Query System' in feature_groups:
        beam_features = [
            "實現 Beam ID 特定的 KPI 查詢",
            "部署專業的 Web UI 界面",
            "完成 E2 Simulator beam_id 支援",
            "實現完整的資料流傳輸",
            "提供 RESTful API 介面",
            "整合系統健康監控報告"
        ]
        add_bullet_slide(prs, "Beam KPI Query System", beam_features)

    # Dual-Path Communication
    if 'Dual-Path Communication' in feature_groups:
        dual_path_features = [
            "實現 RMR + HTTP 雙路徑通訊",
            "整合 Grafana 監控儀表板",
            "為所有 xApps 啟用雙路徑",
            "實現 Web UI 代理伺服器",
            "完成 WSGI 應用程式部署",
            "建立完整的監控架構"
        ]
        add_bullet_slide(prs, "Dual-Path Communication", dual_path_features)

    # Monitoring & Visualization
    if 'Monitoring & Visualization' in feature_groups:
        monitoring_features = [
            "部署 Prometheus 監控系統",
            "建立 Grafana 儀表板",
            "為所有 xApps 添加 HTTP endpoints",
            "實現健康檢查探針",
            "自動化儀表板匯入",
            "完成性能測試框架"
        ]
        add_bullet_slide(prs, "監控與視覺化", monitoring_features)

    # Platform & Infrastructure
    platform_features = [
        "完成 K3s 叢集部署腳本",
        "實現 KUBECONFIG 標準化",
        "建立一鍵部署流程",
        "整合 GPU 支援",
        "部署 E2 Simulator",
        "完成平台 Phase 3 & 4"
    ]
    add_bullet_slide(prs, "平台與基礎設施", platform_features)

    # Documentation Section
    add_section_header(prs, "文檔與測試")

    doc_items = [
        "重組文檔結構",
        "專業化 README",
        "完整的部署指南",
        "故障排除文檔",
        "API 使用指南",
        "快速入門教學"
    ]

    test_items = [
        "單元測試覆蓋",
        "整合測試實作",
        "端到端測試",
        "性能測試報告",
        "E2E 測試文檔",
        "TDD 開發流程"
    ]

    add_two_column_slide(prs, "文檔與測試", doc_items, test_items)

    # Code Quality Section
    add_section_header(prs, "程式碼品質提升")

    quality_bullets = [
        "移除硬編碼路徑，改用動態解析",
        "標準化 Shell 腳本執行權限",
        "清理重複與過時檔案",
        "實現專業的錯誤處理",
        "重構程式碼結構提升可讀性",
        "遵循 PEP 8 與最佳實踐"
    ]
    add_bullet_slide(prs, "程式碼品質改進", quality_bullets)

    # Bug Fixes & Improvements
    fix_bullets = [
        "修復 Docker 映像版本問題",
        "解決 K3s Cilium 超時問題",
        "修正 iptables 清理流程",
        "更新 NGINX 版本相容性",
        "修復 ImagePullBackOff 錯誤",
        "解決 Redis 與 ricsdl 相依性"
    ]
    add_bullet_slide(prs, "錯誤修復", fix_bullets)

    # Recent Progress (Last 7 days)
    add_section_header(prs, "近期進度 (最近一週)")

    recent_commits = [c for c in commits if c['date'] >= '2025-11-15']
    recent_bullets = []

    for commit in recent_commits[:10]:
        msg = commit['message']
        if len(msg) > 70:
            msg = msg[:67] + "..."
        recent_bullets.append(f"{commit['date']}: {msg}")

    add_bullet_slide(prs, "最近一週提交記錄", recent_bullets)

    # Technical Architecture
    add_section_header(prs, "技術架構")

    left_items = [
        "RIC Platform (Near-RT RIC)",
        "xApp Framework",
        "- KPIMON xApp",
        "- RC (RAN Control) xApp",
        "- Traffic Steering xApp",
        "- UAV Policy xApp",
        "E2 Interface",
        "- E2 Simulator",
        "- E2 Termination"
    ]

    right_items = [
        "Monitoring Stack",
        "- Prometheus",
        "- Grafana Dashboards",
        "Communication",
        "- RMR (RIC Message Router)",
        "- HTTP REST API",
        "Infrastructure",
        "- K3s Kubernetes",
        "- Docker Containers"
    ]

    add_two_column_slide(prs, "系統架構組件", left_items, right_items)

    # Current Status
    add_section_header(prs, "目前狀態")

    status_bullets = [
        "✓ 核心平台部署完成",
        "✓ 主要 xApps 正常運作",
        "✓ 監控系統完整建立",
        "✓ UAV Policy xApp 整合完成",
        "✓ Beam KPI 查詢系統上線",
        "✓ 雙路徑通訊機制運作中",
        "✓ 完整的測試覆蓋",
        "✓ 文檔與部署指南齊全"
    ]
    add_bullet_slide(prs, "專案完成度", status_bullets)

    # Future Work
    add_section_header(prs, "未來規劃")

    future_bullets = [
        "持續優化系統性能",
        "擴展監控指標與告警",
        "增強 xApp 功能與自動化",
        "改進測試覆蓋率",
        "整合更多 RAN 場景",
        "生產環境部署準備",
        "安全性強化",
        "文檔持續更新"
    ]
    add_bullet_slide(prs, "下一步工作", future_bullets)

    # Thank You Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Add text box for "Thank You"
    left = Inches(1)
    top = Inches(3)
    width = Inches(8)
    height = Inches(2)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "謝謝聆聽"

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)

    # Add subtitle
    left = Inches(1)
    top = Inches(4.5)

    txBox2 = slide.shapes.add_textbox(left, top, width, Inches(1))
    tf2 = txBox2.text_frame
    tf2.text = "O-RAN RIC Platform 開發團隊"

    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    p2.font.size = Pt(24)
    p2.font.color.rgb = RGBColor(100, 100, 100)

    # Save presentation
    output_file = 'ORAN_RIC_Progress_Report.pptx'
    prs.save(output_file)
    print(f"\n✓ Presentation saved as: {output_file}")
    print(f"  Total slides: {len(prs.slides)}")

    return output_file

if __name__ == "__main__":
    print("=" * 60)
    print("O-RAN RIC Platform - Progress Report Generator")
    print("=" * 60)

    try:
        output_file = generate_presentation()
        print("\n" + "=" * 60)
        print("SUCCESS: PowerPoint presentation generated!")
        print(f"File: {output_file}")
        print("=" * 60)
    except Exception as e:
        print(f"\nERROR: {e}")
        import traceback
        traceback.print_exc()
